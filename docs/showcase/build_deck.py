"""Build the ptm-ipf showcase deck (PowerPoint) from the campaign figures.

Every picture on a slide is a file under ~/showcase/figures, produced by
showcase_run.py or compare_textures.py. Nothing here recomputes physics: the
numbers come from stats.json, peaks.txt and the campaign README.

    python build_deck.py [out.pptx]
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from PIL import Image, ImageDraw, ImageFont

HOME = Path.home()
FIG = HOME / "showcase" / "figures"
OUT = Path(sys.argv[1]) if len(sys.argv) > 1 else HOME / "showcase" / "report" / "ptm-ipf_showcase.pptx"

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x1A, 0x1A, 0x1E)
MUTED = RGBColor(0x5A, 0x5F, 0x6A)
ACCENT = RGBColor(0xB0, 0x3A, 0x2B)
RULE = RGBColor(0xD8, 0xD4, 0xCC)
PAPER = RGBColor(0xFA, 0xF8, 0xF5)
BAND = RGBColor(0x22, 0x26, 0x2B)

FONT = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- primitives

def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BAND if dark else PAPER
    return s


def text(s, left, top, width, height, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, font=FONT, space_after=6, line=None):
    box = s.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def para(tf, run_text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
         font=FONT, space_before=0, space_after=6, first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = run_text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def rule(s, left, top, width, color=RULE, thick=Pt(1.5)):
    from pptx.enum.shapes import MSO_SHAPE
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thick)
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


_MPL = Path(__import__("matplotlib").__file__).parent / "mpl-data" / "fonts" / "ttf"
_MEASURE = ImageDraw.Draw(Image.new("RGB", (8, 8)))


def title_lines(title, size_pt, width_in):
    """How many lines the title needs, measured in DejaVu Bold.

    DejaVu is wider than Calibri, so a title that fits here fits in PowerPoint.
    """
    f = ImageFont.truetype(str(_MPL / "DejaVuSans-Bold.ttf"), int(size_pt * 96 / 72))
    limit = width_in * 96
    lines, line = 1, ""
    for word in title.split():
        trial = (line + " " + word).strip()
        if _MEASURE.textlength(trial, font=f) <= limit or not line:
            line = trial
        else:
            lines += 1
            line = word
    return lines


def heading(s, title, kicker=None):
    top = Inches(0.42)
    if kicker:
        _, tf = text(s, Inches(0.6), top, Inches(12.1), Inches(0.3))
        para(tf, kicker.upper(), size=12, bold=True, color=ACCENT, first=True, space_after=2)
        top = Inches(0.72)
    size = 30 if title_lines(title, 30, 12.1) == 1 else 26
    lines = title_lines(title, size, 12.1)
    _, tf = text(s, Inches(0.6), top, Inches(12.1), Inches(0.6))
    para(tf, title, size=size, bold=True, first=True, space_after=0)
    drop = Inches(0.62) + Inches(0.42) * (lines - 1)
    rule(s, Inches(0.6), top + drop, Inches(12.1))
    return top + drop + Inches(0.23)


def fit(path, left, top, max_w, max_h):
    """Return the (left, top, width, height) that centres an image in a box."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return Emu(int(left + (max_w - w) / 2)), Emu(int(top + (max_h - h) / 2)), Emu(w), Emu(h)


def picture(s, path, left, top, max_w, max_h):
    l, t, w, h = fit(Path(path), int(left), int(top), int(max_w), int(max_h))
    return s.shapes.add_picture(str(path), l, t, w, h)


def caption(s, words, left=Inches(0.6), top=Inches(6.95), width=Inches(12.1)):
    _, tf = text(s, left, top, width, Inches(0.4))
    para(tf, words, size=12, color=MUTED, italic=True, first=True, space_after=0)


# ------------------------------------------------------------- slide recipes

def title_slide(title, subtitle, byline):
    s = slide(dark=True)
    n = title_lines(title, 44, 11.3)
    top = Inches(2.5) - Inches(0.31) * n
    _, tf = text(s, Inches(1.0), top, Inches(11.3), Inches(0.6))
    para(tf, title, size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), first=True, space_after=0)
    y = top + Inches(0.62) * n + Inches(0.25)
    _, tf = text(s, Inches(1.0), y, Inches(11.3), Inches(0.8))
    para(tf, subtitle, size=20, color=RGBColor(0xC8, 0xC2, 0xB8), first=True, space_after=0)
    y += Inches(0.34) * title_lines(subtitle, 20, 11.3) + Inches(0.35)
    rule(s, Inches(1.0), y, Inches(2.4), color=ACCENT, thick=Pt(3))
    _, tf = text(s, Inches(1.0), y + Inches(0.32), Inches(11.3), Inches(0.5))
    para(tf, byline, size=15, color=RGBColor(0x9A, 0x96, 0x90), first=True, space_after=0)
    return s


def section_slide(number, title, blurb):
    s = slide(dark=True)
    _, tf = text(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.0))
    para(tf, number, size=16, bold=True, color=ACCENT, first=True, space_after=10)
    para(tf, title, size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), space_after=12)
    para(tf, blurb, size=17, color=RGBColor(0xB8, 0xB4, 0xAE), space_after=0)
    return s


def bullets_slide(title, items, kicker=None, note=None):
    s = slide()
    top = heading(s, title, kicker)
    _, tf = text(s, Inches(0.6), top + Inches(0.2), Inches(12.1), Inches(5.0))
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            lead, rest = item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(13)
            r = p.add_run(); r.text = lead + "  "
            r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = FONT
            r = p.add_run(); r.text = rest
            r.font.size = Pt(19); r.font.color.rgb = INK; r.font.name = FONT
        else:
            para(tf, item, size=19, first=(i == 0), space_after=13)
    if note:
        caption(s, note)
    return s


def figure_slide(title, image, points, kicker=None, note=None, split=0.60):
    """Picture on the left, a short argument on the right."""
    s = slide()
    top = heading(s, title, kicker)
    box_w = int((W - Inches(1.2)) * split)
    body_h = int(Inches(6.85) - top)
    l, t, w, h = fit(Path(image), int(Inches(0.6)), int(top), box_w, body_h)
    # Left align the picture rather than centring it in its box: a narrow figure
    # otherwise leaves a gap on both sides and crowds the text against nothing.
    pic = s.shapes.add_picture(str(image), Inches(0.6), t, w, h)
    tx = pic.left + pic.width + Inches(0.4)
    tw = W - Inches(0.6) - tx
    if tw < Inches(3.2):
        tx, tw = Inches(0.6) + box_w + Inches(0.35), W - Inches(0.95) - box_w - Inches(0.6)
    _, tf = text(s, Emu(int(tx)), top + Inches(0.1), Emu(int(tw)), Inches(5.0))
    for i, item in enumerate(points):
        if isinstance(item, tuple):
            lead, rest = item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(11)
            r = p.add_run(); r.text = lead
            r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = FONT
            r = p.add_run(); r.text = "\n" + rest
            r.font.size = Pt(15); r.font.color.rgb = INK; r.font.name = FONT
        else:
            para(tf, item, size=15, first=(i == 0), space_after=11)
    if note:
        caption(s, note)
    return s


def big_figure_slide(title, image, note, kicker=None):
    s = slide()
    top = heading(s, title, kicker)
    lines = 1 + len(note) // 150
    body = int(Inches(6.9) - Inches(0.22) * lines - top)
    pic = picture(s, image, Inches(0.6), top, int(W - Inches(1.2)), body)
    caption(s, note, top=Emu(pic.top + pic.height + Inches(0.1)))
    return s


def pair_slide(title, left_img, right_img, left_label, right_label, note, kicker=None):
    s = slide()
    top = heading(s, title, kicker)
    half = int((W - Inches(1.5)) / 2)
    body_h = int(Inches(6.9) - Inches(0.22) * (1 + len(note) // 150) - top)
    for i, (img, label) in enumerate(((left_img, left_label), (right_img, right_label))):
        x = Inches(0.6) + i * (half + Inches(0.3))
        _, tf = text(s, Emu(int(x)), top, Emu(half), Inches(0.32))
        para(tf, label, size=15, bold=True, color=INK, align=PP_ALIGN.CENTER, first=True, space_after=0)
        pic = picture(s, img, int(x), int(top + Inches(0.42)), half, body_h - int(Inches(0.42)))
    caption(s, note, top=Emu(pic.top + pic.height + Inches(0.12)))
    return s


def table_slide(title, headers, rows, kicker=None, note=None, widths=None, size=12):
    s = slide()
    top = heading(s, title, kicker)
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = min(Inches(5.9), Inches(0.32) * n_rows)
    shape = s.shapes.add_table(n_rows, n_cols, Inches(0.6), top + Inches(0.1),
                               W - Inches(1.2), height)
    table = shape.table
    table.first_row = True
    if widths:
        total = sum(widths)
        for c, wfrac in enumerate(widths):
            table.columns[c].width = Emu(int((W - Inches(1.2)) * wfrac / total))
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = head
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.runs[0].font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = BAND
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value) if str(value).strip() else "pending"
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(size)
            p.runs[0].font.color.rgb = INK
            p.runs[0].font.name = FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if r % 2 else RGBColor(0xF0, 0xEC, 0xE6)
    if note:
        caption(s, note, top=top + Inches(0.2) + height + Inches(0.28))
    return s


# ------------------------------------------------------------------ the data

PEAKS = {}
for line in (HOME / "showcase" / "peaks.txt").read_text().splitlines():
    m = re.match(r"(\S+) peak ([\d.]+) GPa at strain ([\d.]+), final ([\d.]+)", line)
    if m:
        PEAKS[m.group(1)] = (float(m.group(2)), float(m.group(3)), float(m.group(4)))

PARENT = {"cu": "fcc", "fe": "bcc", "ti": "hcp"}
OTHER = {"cu": "hcp", "fe": "hcp", "ti": "fcc"}


def stats(name):
    p = FIG / name / "stats.json"
    return json.loads(p.read_text()) if p.exists() else None


def row(name):
    d = stats(name)
    el, tex, mode = name.split("_")[0], name.split("_")[1], name.split("_")[-1]
    load = "tension" if mode == "T" else "compression"
    pk = PEAKS.get(name)
    if d is None:
        return [name, "", tex, load, "", "", "", ""]
    n, c = d["n_atoms"], d["final_counts"]
    frac = 100 * c.get(OTHER[el], 0) / n
    g = d["grains_in_sections"]
    return [name, f"{n:,}", tex, load,
            f"{pk[0]:.2f}" if pk else "", f"{100 * pk[1]:.1f}%" if pk else "",
            f"{OTHER[el]} {frac:.1f}%",
            f"{g['x']}, {g['y']}, {g['z']}"]


RUNS = ["cu_random_T", "cu_rolled_T", "cu_rolled_C", "cu_extruded_T", "cu_extruded_C",
        "fe_random_v2_T", "fe_rolled_v2_C", "ti_random_T", "ti_rolled_C", "ti_extruded_T"]


# ---------------------------------------------------------------- the slides

title_slide(
    "Reading a simulated microstructure like an EBSD map",
    "ptm-ipf: inverse pole figure colouring for atomistic simulations, and a ten run "
    "deformation campaign in Cu, Fe and Ti",
    "Alexander Goldman  |  github.com/Carbonpoint/ptm-ipf  |  GPL-3.0  |  August 2026",
)

bullets_slide(
    "The gap this fills",
    [("OVITO", "computes per atom orientations with polyhedral template matching, but colours "
      "them through Rodrigues space. That is a valid colouring and not an inverse pole figure, "
      "so it cannot be compared with an EBSD map."),
     ("MTEX and orix", "implement the proper EDAX/TSL inverse pole figure keys, but they work "
      "on EBSD scans, not on atoms."),
     ("ptm-ipf", "puts the standard key on the atoms. The same colour means the same "
      "crystal direction in a simulation and in a scan, so the two can sit side by side in "
      "one figure and be read the same way."),
     ("Scope", "hcp, fcc, bcc, simple cubic, cubic and hexagonal diamond and graphene; any "
      "reference direction, named or arbitrary; maps, sections, pole figures, orientation "
      "densities, boundary maps, unit cell overlays and animations.")],
    kicker="Motivation",
)

bullets_slide(
    "What it produces",
    [("Coloured atoms", "the configuration is written back with Color.R/G/B columns that OVITO "
      "binds on reload, so the atoms open already coloured."),
     ("Flat maps", "an EBSD style orientation map of a section: no atoms, just colours, "
      "segmented grains and boundaries, at a chosen pixel size."),
     ("Sections and renders", "slabs of chosen thickness normal to x, y or z, orthographic, "
      "with a sample triad."),
     ("Pole figures and IPF densities", "equal area, in multiples of a random distribution, in "
      "the sample frame you define."),
     ("Boundary analysis", "boundaries coloured by grain to grain misorientation on a chosen "
      "scale, or split at a threshold into high and low angle."),
     ("Unit cell overlays", "a hexagonal prism or a cube drawn on each distinct orientation, "
      "sized by grain area."),
     ("Animations", "any of the above through a trajectory, as MP4 or GIF.")],
    kicker="Capabilities",
)

bullets_slide(
    "Verification, before any picture means anything",
    [("Against orix", "the key agrees with orix to better than 2e-5 in RGB for m-3m, 6/mmm, "
      "4/mmm, -3m and mmm."),
     ("PTM conventions pinned by tests", "the quaternion order and the hcp template frame the "
      "colours depend on are established by tests that run real PTM, so an OVITO change cannot "
      "silently rotate a map."),
     ("Round trip through the builder", "every cell was built from explicit rotation matrices "
      "recorded beside it. PTM recovers each built grain to 0.000 degrees."),
     ("Against published work", "run on the author's own magnesium chapter 6 configurations "
      "the tool reproduces the twin area fractions independently, 26.8 to 35.3 percent against "
      "the 27 to 35 percent measured by hand."),
     ("303 tests, green in CI", "on Python 3.11 and 3.13, with the render tests skipped where "
      "no OpenGL runtime exists.")],
    kicker="Validation",
)

section_slide("Part one", "The campaign",
              "Ten molecular dynamics runs made to exercise every feature on materials whose "
              "behaviour is in the literature.")

table_slide(
    "How the cells were built and loaded",
    ["element", "potential", "structure", "textures", "loading", "rate", "strain"],
    [["Cu", "Mishin 2001, eam/alloy", "fcc", "random, rolled (brass + copper), extruded (<111> + <100>)", "tension and compression", "0.001/ps", "15%"],
     ["Fe", "Mendelev 2003, eam/fs", "bcc", "random, rolled (alpha fibre)", "tension and compression", "0.001/ps", "15%"],
     ["Ti", "Mendelev 2016, eam/fs", "hcp", "random, rolled (basal), extruded (<10-10> fibre)", "tension and compression", "0.0005/ps", "12%"]],
    kicker="Method",
    size=13,
    note="Twenty grain Voronoi polycrystals in 200 A periodic boxes, 410k to 690k atoms. Minimised, "
         "equilibrated 20 ps at 300 K and zero pressure, then strained along x at constant engineering "
         "rate with the transverse directions held at zero pressure. Frames every 10 ps. LAMMPS "
         "29 Aug 2024, CPU only, about 1.5 hours per run on 32 cores.",
)

table_slide(
    "The ten runs at a glance",
    ["run", "atoms", "texture", "load", "peak (GPa)", "at strain", "close packed", "grains per section"],
    [row(n) for n in RUNS],
    kicker="Results",
    size=12,
    widths=[15, 9, 9, 12, 9, 8, 13, 18],
    note="Peak stress is the absolute tensile stress along x. The close packed fraction is the atoms PTM "
         "assigns to hcp or fcc at the final frame: in copper that is stacking faults and twins, a real "
         "quantity; in iron and titanium the few percent is thermal misindexing of a distorted parent "
         "lattice, not a phase. Grain counts include the cells that slip produces inside the built grains.",
)

section_slide("Part two", "Copper: twinning and its asymmetry",
              "Face centred cubic on the Mishin 2001 potential. Every texture twins, and the "
              "sign of the tension and compression asymmetry follows the texture.")

figure_slide(
    "Random copper in tension: twins appear as Sigma 3 lamellae",
    FIG / "cu_random_T" / "boundaries.png",
    [("Yellow at 60 degrees", "the Sigma 3 twin relation, drawn as parallel lamellae through "
      "most of the grains."),
     ("Purple loops", "the low angle cells the grains started with, and those slip creates."),
     ("Nothing in between", "the boundary population is bimodal, which is what a twinning fcc "
      "metal should give."),
     ("12.0 percent hcp", "at the final frame: in fcc copper that is stacking faults and twins, "
      "not a phase change."),
     ("279 to 525 grains", "counted per section, from the 20 that were built.")],
    kicker="Copper, random texture, 632k atoms, 15 percent tension",
    note="Flat map of the 10 A mid cell section normal to z, boundaries coloured by grain to grain "
         "misorientation on a 0 to 60 degree plasma scale.",
)

big_figure_slide(
    "The same section, three copper cells, one colour key",
    FIG / "compare_cu_all_ipfx.png",
    "Every panel is the IPF-X projection of a 10 A section through the middle of the cell, boundary "
    "atoms filled by a 6 A neighbour average, unit cells drawn on each distinct orientation and scaled "
    "to grain area. Rows are the three orthogonal sections; columns are random tension, rolled tension "
    "and rolled compression.",
    kicker="Texture and loading",
)

figure_slide(
    "Rolled copper: three times more twinning in tension than in compression",
    FIG / "compare_cu_rolled_loading_ipfx.png",
    [("Same starting cell", "cu_rolled_T and cu_rolled_C are the identical structure, pulled and "
      "pushed."),
     ("Tension", "12.4 percent hcp, 514 grains counted in the z section, most of them twin "
      "lamellae."),
     ("Compression", "8.3 percent hcp and 274 grains: thin lamellae in the blue grain and one "
      "wide deformation band through the pink one."),
     ("Why", "the Schmid factors of the leading and trailing partials swap with the sign of the "
      "load, so a texture that favours twinning one way suppresses it the other.")],
    kicker="Rolled copper, pulled against pushed",
    split=0.55,
)

figure_slide(
    "Extruded copper: the asymmetry runs the other way",
    FIG / "compare_cu_extruded_loading_ipfx.png",
    [("Tension", "11.2 percent hcp, 166 to 261 grains per section."),
     ("Compression", "13.7 percent hcp and 426 to 690 grains, the most fragmented "
      "microstructure of the whole campaign."),
     ("Why the flip", "for a <111> fibre loaded along its axis the leading partial has the "
      "higher Schmid factor in compression, so pushing the fibre twins it; for the rolled "
      "components the same argument favours tension."),
     ("The point for the tool", "same code, same sections, same key, and it separates two "
      "opposite answers that crystal plasticity predicts.")],
    kicker="Extruded copper, pulled against pushed",
    split=0.55,
)

section_slide("Part three", "Titanium: fragmentation without twinning",
              "Hexagonal close packed on the Mendelev 2016 potential. None of the three textures "
              "twins at this grain size and rate; the texture sets how finely the grains break up.")

big_figure_slide(
    "Three titanium textures, same sections, same key",
    FIG / "compare_ti_all_ipfx.png",
    "Random tension, basal rolled compression and <10-10> extruded tension. The boundary population is "
    "ordinary grain boundaries plus low angle cells inside the grains: there are no twin lamellae in any "
    "of the three.",
    kicker="Titanium, all three textures",
)

figure_slide(
    "Basal titanium in compression: loaded the wrong way to twin",
    FIG / "ti_rolled_C" / "pole_figures.png",
    [("A single 180 MRD basal spot", "at ND: the built texture, and it survives 12 percent "
      "compression essentially unchanged."),
     ("No extension twinning", "compression perpendicular to the c axis cannot drive {10-12} "
      "twinning, which needs the c axis to lengthen."),
     ("Least fragmented cell of the set", "30 to 88 grains per section, against 150 to 321 for "
      "random titanium."),
     ("The animation settles it", "the few twin nuclei appear early and never grow, so this is "
      "a cell that did not twin rather than one that twinned and detwinned. A still map cannot "
      "tell those apart.")],
    kicker="Titanium, basal texture, 12 percent compression along RD",
    note="Equal area pole figures in the sample frame, in multiples of a random distribution.",
)

bullets_slide(
    "What the titanium runs mean",
    [("Extruded tension should have twinned", "ti_extruded_T is built with every c axis "
      "perpendicular to the tensile axis, the loading {10-12} extension twinning needs, and the "
      "section still shows a single green colour family with no lamellae."),
     ("Two consistent explanations", "the Mendelev 2016 potential has a high twin nucleation "
      "stress, and with 20 grains in a 200 A box at 0.0005/ps there is little room for a twin to "
      "nucleate before slip relieves the stress."),
     ("So the texture signal is fragmentation, not twinning",
      "30 to 88 cells per section for basal compression, 91 to 220 for the <10-10> fibre in "
      "tension, 150 to 321 for random."),
     ("Reported as measured", "the tool is not told what to look for. It plots the orientation "
      "field, and the absence of a twin population is as much a result as its presence.")],
    kicker="Titanium",
)

section_slide("Part four", "Iron, and what a flat map caught",
              "The one place in the campaign where the tool found a fault in the input rather "
              "than a feature of the metal.")

pair_slide(
    "The same cell, two builders",
    FIG / "fe_random_T" / "boundaries.png",
    FIG / "fe_random_v2_T" / "boundaries.png",
    "v1 builder: a large unindexable region",
    "v2 builder (atomsk): fully indexed",
    "Flat maps of the 10 A mid cell section normal to z, boundaries coloured by misorientation. The white "
    "region on the left is not a phase, a twin or a cutoff effect: it is present at the first frame before "
    "any strain, and no fill radius closes it.",
    kicker="Random iron, first build against rebuild",
)

bullets_slide(
    "Diagnosing the white region",
    [("What it was not", "not a phase: PTM assigns it nothing at all. Not strain induced: it is "
      "there in frame zero. Not the PTM RMSD cutoff: it survives every value tried, and no "
      "boundary fill radius closes it."),
     ("What it was", "a poorly crystallised zone left by the first builder, which placed rotated "
      "lattice blocks by hand and trimmed the overlaps, leaving the cell at 82 to 93 percent of "
      "ideal density. 20 ps at 300 K cannot recrystallise 30 A of disorder."),
     ("The fix", "build_poly2.py hands the Voronoi geometry to atomsk, which fills it at full "
      "density. The rebuilt cells are 687,259 and 687,273 atoms at 100.0 percent of ideal, "
      "against 610,364 and 643,992 before."),
     ("The convention that had to be established", "atomsk node angles are not Euler angles: "
      "they are extrinsic rotations about x, y and z composing as Rz Ry Rx, pinned down by "
      "building single grains at known angles and measuring them back with PTM."),
     ("The tool was right the whole time", "an unindexable region is the correct report for "
      "atoms that match no template. Reading it as a physical result would have been the error.")],
    kicker="Iron",
)

figure_slide(
    "What the artefact cost in the mechanics",
    FIG / "fe_builder_stress.png",
    [("Random tension", "peak 4.99 GPa on the v1 cell, 5.94 GPa on the rebuilt one, and a "
      "visibly lower elastic slope: the missing 11 percent of atoms softened the whole curve."),
     ("Rolled compression", "5.80 GPa becomes 6.06 GPa. That cell was denser to begin with, at "
      "93.7 percent, and the two curves track each other until yield."),
     ("The lesson", "a density deficit at the grain boundaries reads as a lower modulus and an "
      "early yield, and nothing in the stress curve alone says why. The orientation map did.")],
    kicker="First build against rebuild",
    split=0.58,
)

big_figure_slide(
    "Rebuilt iron: two textures, same sections, same key",
    FIG / "compare_fe_texture_ipfx.png",
    "Random tension and rolled alpha fibre compression, IPF-X, 10 A sections. Iron deforms by slip in both "
    "textures: no twin lamellae anywhere, ordinary 40 to 50 degree boundaries with low angle cells inside "
    "the grains. It is the strongest material of the set by a factor of two.",
    kicker="Rebuilt iron, both textures",
)

section_slide("Part five", "The features, one figure each",
              "Everything above was drawn by the same library. These are the pieces.")

pair_slide(
    "Boundary analysis: by misorientation, or split at a threshold",
    FIG / "cu_extruded_C" / "boundaries.png",
    FIG / "cu_extruded_C" / "hagb_lagb.png",
    "coloured by grain to grain misorientation, 0 to 60 degrees",
    "high angle black, low angle grey, threshold at 5 degrees",
    "Grains are segmented on the map, a mean orientation is taken for each with the symmetry aware "
    "average, and every boundary segment is coloured by the disorientation of the pair it separates. "
    "The axis, the range and the colormap are all arguments.",
    kicker="Features",
)

figure_slide(
    "Pole figures and orientation densities, in the sample frame you define",
    FIG / "fe_rolled_v2_C" / "pole_figures.png",
    [("Equal area, in MRD", "Lambert projection, contoured in multiples of a random "
      "distribution, so the numbers mean what they mean in a texture paper."),
     ("Your sample frame", "RD, TD, ND or ED are whatever vectors you say they are; the "
      "projection follows."),
     ("Reads the built texture back", "the alpha fibre appears as a single sharp {110} spot at "
      "RD, which is how the cell was constructed."),
     ("Also available", "the orientation density inside the fundamental sector, on the same "
      "colour key as the maps.")],
    kicker="Features",
    split=0.56,
)

big_figure_slide(
    "Sections with unit cell overlays",
    FIG / "cu_rolled_C" / "panels" / "cu_rolled_C_cutz_ipfx.png",
    "A 10 A section normal to z, IPF-X, boundary atoms filled by a 6 A neighbour average, with a cube "
    "drawn on every distinct orientation, scaled to grain area and coloured as the inverse of what is "
    "underneath. Cubes in every attitude mean random; aligned pairs mean rolled components.",
    kicker="Features",
)

pair_slide(
    "Atoms, not just maps",
    FIG / "fe_random_v2_T" / "final_ipf_x.png",
    FIG / "fe_random_v2_T" / "panels" / "fe_random_v2_T_render_cutz_ipfx.png",
    "the full cell, cut open, IPF-X on the atoms",
    "the same section rendered orthographically",
    "The colours are written back onto the atoms as Color.R/G/B columns, so the configuration opens in "
    "OVITO already coloured. The sample triad is placed by probing the rendered image until the corner "
    "is clear of the scene.",
    kicker="Features",
)


def video_slide(title, video, poster, points, kicker=None, note=None):
    s = slide()
    top = heading(s, title, kicker)
    box_w = int((W - Inches(1.2)) * 0.56)
    body_h = int(Inches(6.85) - top)
    l, t, w, h = fit(Path(poster), int(Inches(0.6)), int(top), box_w, body_h)
    s.shapes.add_movie(str(video), l, t, w, h, poster_frame_image=str(poster),
                       mime_type="video/mp4")
    tx = Inches(0.6) + box_w + Inches(0.35)
    _, tf = text(s, tx, top + Inches(0.1), W - tx - Inches(0.6), Inches(5.0))
    for i, item in enumerate(points):
        lead, rest = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        r = p.add_run(); r.text = lead
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = FONT
        r = p.add_run(); r.text = "\n" + rest
        r.font.size = Pt(15); r.font.color.rgb = INK; r.font.name = FONT
    if note:
        caption(s, note)
    return s


video_slide(
    "Animations: the map through the whole trajectory",
    FIG / "cu_random_T" / "evolution.mp4",
    FIG / "cu_random_T" / "evolution_mid.png",
    [("Plays in the deck", "click it in slide show. The same frames are in the repository as MP4 "
      "and GIF."),
     ("What it adds", "twins appear one at a time from about 4 percent strain. A single final "
      "frame cannot tell a cell that never twinned from one that twinned and detwinned; the "
      "animation can, and that is how the basal titanium result was settled."),
     ("Frames are padded, not resized", "a deforming cell changes aspect from frame to frame, so "
      "each still is padded on white to a common even size and the scale bar stays honest.")],
    kicker="Features",
    note="cu_random_T, flat map of the z section, 16 frames over 15 percent tension, boundaries coloured "
         "by misorientation.",
)

section_slide("Part six", "What this cost, and what to take from it",
              "Ten runs, one library, and a short list of things that were learned the hard way.")

bullets_slide(
    "Lessons from the campaign",
    [("One PTM cutoff does not fit every lattice", "the RMSD cutoff of 0.1 that suits fcc copper "
      "is too tight for bcc iron at 300 K, whose larger thermal displacements push about 3 "
      "percent of grain interiors over it. Both iron figure sets use 0.15."),
     ("Build at full density or do not trust the mechanics", "the hand rolled builder cost about "
      "1 GPa of peak stress and a visibly lower modulus in random iron, and produced a region "
      "the analysis could not index at all. atomsk fills the same geometry at 100.0 percent."),
     ("Check the conventions, do not assume them", "PTM quaternions, the hcp template frame and "
      "atomsk node angles were each established by measurement rather than from documentation, "
      "and two of the three were not what the obvious reading suggested."),
     ("An unindexable region is a result", "the correct behaviour for a tool is to report that "
      "it cannot index something, and the correct response is to ask why, not to widen the "
      "cutoff until it goes away.")],
    kicker="Method",
)

bullets_slide(
    "Cost and reproducibility",
    [("Molecular dynamics", "about 1.5 hours per run on 32 CPU cores, 12 million atom steps per "
      "second. GPU acceleration was abandoned: the LAMMPS GPU package would not re initialise "
      "between run stages under CUDA 13."),
     ("Analysis", "6 to 7 minutes for a full figure set per run, with the flat map animations at "
      "one frame every 5 to 10 seconds."),
     ("Environments", "uv builds the whole analysis environment in about 8 seconds, and the same "
      "lockstep environment runs in CI on Python 3.11 and 3.13."),
     ("Everything is in the repository", "the builder, the run driver, the comparison figure "
      "script, the summary table generator and the campaign write up, beside the figures they "
      "produce.")],
    kicker="Practicalities",
)

s = section_slide("", "github.com/Carbonpoint/ptm-ipf",
                  "GPL-3.0. Install with uv, run from the command line or the web interface, or "
                  "import the library. Issues and pull requests welcome.")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"{OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
      f"{OUT.stat().st_size / 1e6:.1f} MB)")

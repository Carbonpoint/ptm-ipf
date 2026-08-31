"""Render a .pptx to PNG previews with PIL, so the layout can be checked by eye.

Not a faithful renderer: it draws the shape boxes, the pictures and the text at
the right positions and sizes, which is enough to catch overflow, collisions and
pictures that do not fit.
"""
import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

SRC = Path(sys.argv[1])
DST = Path(sys.argv[2]) if len(sys.argv) > 2 else SRC.parent / "preview"
DST.mkdir(parents=True, exist_ok=True)
DPI = 96
MPL = Path(__import__("matplotlib").__file__).parent / "mpl-data" / "fonts" / "ttf"
FONTS = {
    (False, False): MPL / "DejaVuSans.ttf",
    (True, False): MPL / "DejaVuSans-Bold.ttf",
    (False, True): MPL / "DejaVuSans-Oblique.ttf",
    (True, True): MPL / "DejaVuSans-BoldOblique.ttf",
}


def px(emu):
    return int(Emu(emu).inches * DPI)


def font(size_pt, bold=False, italic=False):
    return ImageFont.truetype(str(FONTS[(bool(bold), bool(italic))]), max(6, int(size_pt * DPI / 72)))


def wrap(draw, words, f, width):
    lines, line = [], ""
    for word in words.split():
        trial = (line + " " + word).strip()
        if draw.textlength(trial, font=f) <= width or not line:
            line = trial
        else:
            lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


prs = Presentation(str(SRC))
W, H = px(prs.slide_width), px(prs.slide_height)
overflows = []

for i, slide in enumerate(prs.slides, start=1):
    bg = (250, 248, 245)
    try:
        c = slide.background.fill.fore_color.rgb
        bg = (c[0], c[1], c[2])
    except Exception:
        pass
    img = Image.new("RGB", (W, H), bg)
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        L, T, SW, SH = px(shape.left), px(shape.top), px(shape.width), px(shape.height)
        if shape.shape_type == 13 or shape.__class__.__name__ == "Picture":
            try:
                pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                img.paste(pic.resize((max(1, SW), max(1, SH))), (L, T))
            except Exception:
                draw.rectangle([L, T, L + SW, T + SH], outline=(200, 0, 0))
            continue
        if shape.has_table:
            table = shape.table
            widths = [px(c.width) for c in table.columns]
            y = T
            for r, trow in enumerate(table.rows):
                rh = px(trow.height)
                x = L
                for c, cell in enumerate(trow.cells):
                    fill = (34, 38, 43) if r == 0 else ((250, 248, 245) if r % 2 else (240, 236, 230))
                    draw.rectangle([x, y, x + widths[c], y + rh], fill=fill, outline=(215, 210, 200))
                    p = cell.text_frame.paragraphs[0]
                    run = p.runs[0] if p.runs else None
                    size = run.font.size.pt if run and run.font.size else 12
                    f = font(size, bold=bool(run and run.font.bold))
                    col = (255, 255, 255) if r == 0 else (26, 26, 30)
                    lines = wrap(draw, cell.text, f, widths[c] - 8)
                    ly = y + 3
                    for line in lines:
                        draw.text((x + 4, ly), line, font=f, fill=col)
                        ly += f.size + 2
                    if ly > y + rh + 4:
                        overflows.append(f"slide {i}: table cell r{r}c{c} overflows its row")
                    x += widths[c]
                y += rh
            continue
        if not shape.has_text_frame:
            try:
                c = shape.fill.fore_color.rgb
                draw.rectangle([L, T, L + SW, T + max(1, SH)], fill=(c[0], c[1], c[2]))
            except Exception:
                pass
            continue
        if not shape.text_frame.text.strip():
            # An autoshape with no text: the thin rules under the headings.
            try:
                c = shape.fill.fore_color.rgb
                draw.rectangle([L, T, L + SW, T + max(1, SH)], fill=(c[0], c[1], c[2]))
            except Exception:
                pass
            continue
        y = T
        for p in shape.text_frame.paragraphs:
            if not p.runs:
                continue
            y += (p.space_before.pt * DPI / 72) if p.space_before else 0
            # Flow the paragraph word by word: a paragraph can mix runs (a bold
            # lead, then body text) and can carry explicit newlines.
            words = []
            for run in p.runs:
                f = font(run.font.size.pt if run.font.size else 18,
                         bold=run.font.bold, italic=run.font.italic)
                col = run.font.color.rgb if run.font.color and run.font.color.type is not None else None
                col = (col[0], col[1], col[2]) if col else (26, 26, 30)
                for j, chunk in enumerate(run.text.split("\n")):
                    if j:
                        words.append(("\n", f, col))
                    for w in chunk.split(" "):
                        if w:
                            words.append((w, f, col))
            if not words:
                continue
            x = L
            line_h = max(f.size for _, f, _ in words) + 5
            for w, f, col in words:
                if w == "\n":
                    x, y = L, y + line_h
                    continue
                adv = draw.textlength(w + " ", font=f)
                if x > L and x + draw.textlength(w, font=f) > L + SW:
                    x, y = L, y + line_h
                draw.text((x, y), w, font=f, fill=col)
                x += adv
            y += line_h + ((p.space_after.pt * DPI / 72) if p.space_after else 0)
        if y > H:
            overflows.append(f"slide {i}: text box at y={T} runs {y - H}px past the bottom")
    img.save(DST / f"{i:02d}.png")

print(f"{len(prs.slides._sldIdLst)} slides -> {DST}")
for line in overflows:
    print("OVERFLOW", line)
